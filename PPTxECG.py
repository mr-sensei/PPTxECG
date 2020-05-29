#!/usr/bin/env python3
'''This tools aims to give guidance toward making better presentations and avoiding death by PowerPoint.
Less is more.
'''
from pptx import Presentation
from xlsxwriter import workbook
from pathlib import Path
import re
import settings
import logging
logging.basicConfig(level=logging.DEBUG, format=' %(levelname)s - %(asctime)s: %(message)s')

#Both main tuples for passing data have the following parts. This is defined here for readability.
PRES_COUNT = 0  # Number of presentations
PRES_DICT = 1   # Data
PRES_PATH = 2   # Folder containing to presentations

#External functions    
def analyse_this(untrusted_path):
    '''Takes path to presentation or set of presentation. Returns tuple packed with raw data.'''
    pres_data = _path_to_pres(untrusted_path)
    if pres_data[0] == 0:
        return pres_data
    metrics = _collect_presentation_metrics(pres_data)
    return metrics

def make_spreadsheet_of_this(my_presentations, hours = 0):
    '''Takes either path to presentation [set] or pre-processed tuple. Creates spreadsheet of data.
    
    analyse_this(path) returns a tuple that this can turn into a spreadsheet (avoids repeating data extraction).
    Optionally accepts hours: if no hours provided, WPM metrics are not shown.'''
    if type(my_presentations) is tuple: #pre-processed tuple?
        mtrcs = my_presentations
    else: #untrusted path to presentation?
        mtrcs = analyse_this(my_presentations)
    
    if mtrcs[0] == 0:
        logging.debug('No data found.')
        pass #in case no metrics returned (e.g. invalid path)
    else:
        try:
            #Make a spreadsheet
            xlsx_path = _make_xlsx(mtrcs, int(hours))
            if xlsx_path:
                return xlsx_path
        except:
            return False
    return False

#Internal Functions
def _collect_presentation_metrics(pres_tuple):
    if(pres_tuple[PRES_COUNT] == 0):
        return None
    else:
        presentation_set = {}
        word_count = {}
        logging.debug('Collecting metrics.')
        for prsnt in pres_tuple[PRES_DICT]:
            logging.debug('Looking at %s (%s)',prsnt, pres_tuple[PRES_DICT][prsnt])
            raw = _get_raw_text(pres_tuple[PRES_DICT][prsnt])
            logging.debug('Requesting word count for %s.',prsnt)
            word_count[prsnt] = _get_word_count(raw)
            logging.debug('Word count returned for %s is %s.',prsnt, word_count[prsnt])
        presentation_set = _merge_data(pres_tuple,word_count)
        return(presentation_set)

def _find_pptx(folder_path):
    logging.debug('Looking for presentations in %s', folder_path)
    pres_listing = []
    for files in sorted(folder_path.glob('*.pptx')):
        pres_listing.append(files)
    if pres_listing:
        logging.debug('Found %s pptx file(s): %s', len(pres_listing), pres_listing)
    else:
        return (0, None)
    return(len(pres_listing), pres_listing)


def _get_slide_comment(sld_wc):
    logging.debug("Seeking comment on wordcount of %s.", sld_wc)
    comment = ''
    for each_count in settings.Word_Counts_Comments:
        if sld_wc <= each_count[0]:
            comment = each_count[1]
            break
    if not comment:
        comment = "Extreme. Off the charts."
    logging.debug("Returning '%s'.",comment)
    return comment


def _get_raw_text(grt_pres_obj):
    logging.debug('Scraping raw text from %s.',grt_pres_obj)
    this_slide = 0
    all_content = {}
    for slide in grt_pres_obj.slides:
        this_slide += 1
        this_content = []
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1: #get title
                if(shape.text_frame.paragraphs[0].runs):
                    this_content.append(shape.text_frame.paragraphs[0].runs[0].text)
            if shape.has_text_frame:        #get text in shapes
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        this_content.append(run.text)
            if shape.has_table:             #get text in tables
                cells = shape.table.iter_cells()
                for each_cell in cells:
                    for paragraph in each_cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            this_content.append(run.text)
        all_content[this_slide] = this_content
    logging.debug('Returning raw text.')
    return all_content


def _get_titles(prs_obj):
    titles = {}
    slide_no = 0
    for slide in prs_obj.slides:
        slide_no += 1
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1: #is a title
                if(shape.text_frame.paragraphs[0].runs):
                    titles[slide_no] = shape.text_frame.paragraphs[0].runs[0].text
        if slide_no not in titles:
            titles[slide_no] = 'No title.'
    logging.debug('Found presentation titles: %s',titles)
    return titles


def _get_word_count(raw_text):
    total_wc = 0
    slide_counts = {}
    for slides in raw_text:
        raw = str(raw_text[slides])
        words = re.findall("[a-zA-Z0-9_]+", raw)
        slide_counts[slides] = len(words)
        total_wc += len(words)
    logging.debug('Returning word count.')
    return (total_wc, slide_counts)

def _get_wpm_comment(all_words, all_hours):
    wpm = int(all_words/(all_hours * 60))
    wpm_comment = ''
    for wpm_counts in range(len(settings.WPM)):
        if wpm > settings.WPM[wpm_counts][0]:
            wpm_comment = settings.WPM[wpm_counts][1]
            break
        wpm_comment = "WPM Excellent. Well designed for a prepared trainer. Make sure the trainer guide is good."
    return wpm, wpm_comment


def _make_xlsx(metrics, hours):
    logging.debug('Creating spreadsheet.')
    ext = '_metrics.xlsx'
    save_path = metrics[PRES_PATH]
    if (metrics[PRES_COUNT] == 1):
        for key in metrics[PRES_DICT]:
            pres_name = str(key)
            file_name = pres_name[0:-5]
    else:
        pres_name = ''
        file_name = save_path.parts[-1]
    file_name += ext
    full_name = Path.joinpath(save_path,file_name)
    logging.debug('Spreadsheet \'%s\' will be created at %s.',file_name,save_path)
    
    #Create the spreadsheet
    spreadsheet = workbook.Workbook(full_name)
    default_cells = spreadsheet.add_format({'font_color':'black'})
    yellow_cells = spreadsheet.add_format({'font_color':'#BBBB00'})
    red_cells = spreadsheet.add_format({'font_color':'red'})
    pres_data = metrics[PRES_DICT]

    #Cover-page
    ################
    # Folder path, no of pres || File name
    #
    # Tot wordcount: ...
    # [ Planned hours: ...
    #   WPM Note: ...       ]
    #
    # Pres name | slide count   | word count    | wps
    # Pres1.pptx| Num           | NUM           | Num
    # ...
    ################
    coversheet = spreadsheet.add_worksheet("Summary")
    row = 0
    col = 0
    #Title
    if (metrics[PRES_COUNT] == 1): #If there is only one presentation, we'll put it here.
        coversheet.write(row, col,str(Path.joinpath(save_path,pres_name)))
        col0_max = len(str(Path.joinpath(save_path,pres_name)))
    else: # If > 1, put folder here and number of presentations.
        coversheet.write(row, col + 0, str(save_path))
        coversheet.write(row, col + 1, str(metrics[PRES_COUNT]) + ' presentations')
        col0_max = len(str(save_path))
    
    col13_max = 14
    row += 2
    
    #Brief
    tot_words = 0
    for key in pres_data:
        tot_words += pres_data[key][0]
    
    coversheet.write(row, col + 0, "Total Wordcount: " + str(tot_words))
    if hours > 0: #Only do this is time specified
        tot_wpm, tot_wpmcomment = _get_wpm_comment(tot_words, hours)
        row += 1
        coversheet.write(row, col + 0, "Planned hours: " + str(hours))
        row += 1
        coversheet.write(row, col + 0, "Words/Min: " + str(tot_wpm) + "WPM - " + tot_wpmcomment)
    row += 2

    #Body
    coversheet.write(row, col + 0, "Title")
    coversheet.write(row, col + 1, "Slide count")
    coversheet.write(row, col + 2, "Word count")
    coversheet.write(row, col + 3, "Words/Slide")

    for key in pres_data:
        row += 1
        coversheet.write(row, col + 0, key)
        coversheet.write(row, col + 1, len(pres_data[key][1]))
        coversheet.write(row, col + 2, pres_data[key][0])
        coversheet.write(row, col + 3, int(pres_data[key][0]/len(pres_data[key][1])))
        if len(key) > col0_max: col0_max = len(key) + 2
    
    coversheet.set_column(0, 0, col0_max)
    coversheet.set_column(1, 3, col13_max)

    logging.debug('Coversheet prepared.')
    # Data-sheets
    ##########
    # Pres Name
    # Slide No  | Slide Title   | Word Count
    # 1         | Blah          | ...
    # ...
    ####
    # Pres Name
    # Slide No  | Slide Title   | Word Count
    # 1         | Blah          | ...
    # ...
    ####
    # ...
    #########
    for keys in pres_data:
        logging.debug('Creating worksheet for %s (%s slides).',keys, pres_data[keys])
        datasheet = spreadsheet.add_worksheet(keys.split()[0][0:30]) #Doesn't allow more than 31 chars - this takes first word and cuts if still too long
        slide_title = _get_titles(pres_data[keys][2])
        col = 0
        row = 0
        # Title
        datasheet.write(row, col + 0, keys)
        row += 1
        # Data
        datasheet.write(row, col + 0, "Slide No.")
        datasheet.write(row, col + 1, "Word count")
        datasheet.write(row, col + 2, "Title")
        datasheet.write(row, col + 3, "Comment")
        col2_max = 10
        col3_max = 12

        for slides in pres_data[keys][1]:
            cell_colour = default_cells
            row += 1
            slide_wc = pres_data[keys][1][slides]
            slide_comment = _get_slide_comment(slide_wc)
            if 'warning' in slide_comment.lower() or 'caution' in slide_comment.lower():
                cell_colour = yellow_cells
            elif 'extreme' in slide_comment.lower():
                cell_colour = red_cells
            datasheet.write(row, col + 0, slides, cell_colour)
            datasheet.write(row, col + 1, slide_wc, cell_colour)
            datasheet.write(row, col + 2, slide_title[slides], cell_colour)
            datasheet.write(row, col + 3, slide_comment, cell_colour)
            if len(slide_title[slides]) > col2_max: col2_max = len(slide_title[slides])
            if len(slide_comment) > col3_max: col3_max = len(slide_comment)

        datasheet.set_column(0, 1, 10) # Col A, B set to 10
        datasheet.set_column(2, 2, col2_max) # Col C dynamic with slide title
        datasheet.set_column(3, 3, col3_max) # Col D dynamic with comment

    spreadsheet.close()   
    return full_name


def _merge_data(p_tuple, mdata_dict):
    '''Takes file tuple and wordcount mdata dict - produces merged tuple.
    ( No_of_presentations, { 'Filename1': (pres_word_count, {slide1: slide_count, s2:..., ...}, <presentation_object>), "File2":(...), ... }, file/folder path ) '''
    merged_dict = {}
    for files in p_tuple[PRES_DICT]:
        merged_dict[files] = (mdata_dict[files][0],mdata_dict[files][1],p_tuple[1][files])
    return(p_tuple[PRES_COUNT],merged_dict,p_tuple[PRES_PATH])


def _path_to_pres(full_path):
    '''Pass in a path to a presentation file or a folder containing presentations.
    Returns tuple containing:

    #### count_of_presentations_found   ###   dict(name:presentation_objects)  ###   file/folder path ####'''
    logging.debug('Passed a path: %s.',full_path)
    pres_path = Path(full_path)
    pres = (0, None, None)
    if pres_path.is_file(): #If it's a file, try to open a presentation object and return it.
        logging.debug('This looks like the file %s.',pres_path.name)
        try:
            pres = (1, {pres_path.name:Presentation(pres_path)}, pres_path.parent)
            logging.debug('Successfully opened presentation.')
        except:
            logging.debug('Failed to open file. This does not look like a valid presentation.')
            pass
    elif pres_path.is_dir(): #If it's a folder, call _find_pptx
        logging.debug('This looks like the directory %s',pres_path)
        count, pres_list = _find_pptx(pres_path)
        if count == 0:
            logging.debug('No presentations found here.')
            pass
        else: #If more than one, make a lookup table of name:pres_obj
            pres_obj_dict = {}
            for prezies in pres_list:
                try:
                    pres_obj_dict[prezies.name] = Presentation(prezies)
                except:
                    logging.debug('Failed to open file %s - it does not look like a valid presentation.',prezies.name)
            if pres_obj_dict:
                pres = (len(pres_obj_dict),pres_obj_dict,pres_path)
                logging.debug('Successfully opened %s object(s).',len(pres_obj_dict))
            else:
                pass
    else:
        logging.debug('Object %s not found.',pres_path)
    logging.debug('PPTxECG will return %s.',pres)
    return pres



if __name__ == "__main__":
    #Switch things to test different scenario
    thing = Path('.','sample/sample.pptx')
    x = analyse_this(thing)
    print(x)
    y = make_spreadsheet_of_this(x)
    print(y)